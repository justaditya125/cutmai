"""
Advanced File Processing Service
Extends the existing FileHandler with:
- PPT/PPTX, SVG, JSON, XML, ZIP/RAR/7Z parsing
- Batch processing
- Duplicate detection (SHA-256)
- File chunking for RAG
"""
import hashlib
import os
from pathlib import Path
from typing import Dict, List, Optional, Tuple


class FileProcessor:
    """Wraps and extends the existing FileHandler with advanced parsing capabilities."""

    def process(self, file_bytes: bytes, filename: str) -> Dict:
        """
        Process a file and return extracted text + metadata.
        Delegates to existing FileHandler for already-supported types,
        and uses new parsers for extended types.
        """
        ext = Path(filename).suffix.lower().lstrip('.')
        text = ''
        metadata = MetadataExtractor.extract(filename, file_bytes)

        try:
            if ext in ('pptx', 'ppt'):
                text = PPTParser.parse(file_bytes)
            elif ext == 'svg':
                text = SVGParser.parse(file_bytes)
            elif ext == 'json':
                text = JSONParser.parse(file_bytes)
            elif ext in ('xml',):
                text = XMLParser.parse(file_bytes)
            elif ext in ('zip', 'rar', '7z', 'tar', 'gz', 'bz2'):
                text = ArchiveParser.list_contents(file_bytes, filename)
            elif ext in ('md', 'markdown', 'txt', 'log', 'rst', 'csv'):
                text = file_bytes.decode('utf-8', errors='replace')
            else:
                # Let existing client-side parsers handle PDF, DOCX, XLSX
                text = f'[File: {filename} — client-side parsing required]'
        except Exception as e:
            print(f"[FileProcessor] Parse error for {filename}: {e}")
            text = f'[Parse error: {e}]'

        return {
            'filename':  filename,
            'extension': ext,
            'text':      text,
            'metadata':  metadata,
            'char_count': len(text)
        }

    def batch_process(self, files: List[Tuple[bytes, str]]) -> List[Dict]:
        """Process multiple files and return results list."""
        results = []
        for file_bytes, filename in files:
            result = self.process(file_bytes, filename)
            results.append(result)
        return results


class PPTParser:
    @staticmethod
    def parse(file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            texts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    texts.append(f'[Slide {slide_num}]\n' + '\n'.join(slide_texts))
            return '\n\n'.join(texts)
        except ImportError:
            return '[PPTX parsing requires python-pptx. Run: pip install python-pptx]'
        except Exception as e:
            return f'[PPTX parse error: {e}]'


class SVGParser:
    @staticmethod
    def parse(file_bytes: bytes) -> str:
        try:
            import re
            svg_text = file_bytes.decode('utf-8', errors='replace')
            # Extract text elements
            text_matches = re.findall(r'<text[^>]*>(.*?)</text>', svg_text, re.DOTALL)
            title_matches = re.findall(r'<title>(.*?)</title>', svg_text, re.DOTALL)
            desc_matches = re.findall(r'<desc>(.*?)</desc>', svg_text, re.DOTALL)
            parts = title_matches + desc_matches + text_matches
            return '\n'.join(p.strip() for p in parts if p.strip()) or '[SVG: no readable text found]'
        except Exception as e:
            return f'[SVG parse error: {e}]'


class JSONParser:
    @staticmethod
    def parse(file_bytes: bytes) -> str:
        import json
        try:
            data = json.loads(file_bytes.decode('utf-8', errors='replace'))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            return f'[JSON parse error: {e}]'


class XMLParser:
    @staticmethod
    def parse(file_bytes: bytes) -> str:
        try:
            try:
                import defusedxml.ElementTree as ET
            except ImportError:
                import xml.etree.ElementTree as ET
            root = ET.fromstring(file_bytes.decode('utf-8', errors='replace'))
            texts = [el.text for el in root.iter() if el.text and el.text.strip()]
            return '\n'.join(texts) or '[XML: no text content]'
        except Exception as e:
            return f'[XML parse error: {e}]'


class ArchiveParser:
    @staticmethod
    def list_contents(file_bytes: bytes, filename: str) -> str:
        import io, zipfile
        try:
            if filename.lower().endswith('.zip'):
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    names = zf.namelist()
                    return f'[ZIP Archive — {len(names)} files]\n' + '\n'.join(names[:50])
            return f'[Archive: {filename} — contents not extracted (non-ZIP format)]'
        except Exception as e:
            return f'[Archive parse error: {e}]'


class MetadataExtractor:
    @staticmethod
    def extract(filename: str, file_bytes: bytes) -> Dict:
        sha256 = hashlib.sha256(file_bytes).hexdigest()
        return {
            'filename':    filename,
            'extension':   Path(filename).suffix.lower(),
            'size_bytes':  len(file_bytes),
            'sha256':      sha256,
        }


class DuplicateDetector:
    """Detects duplicate files using SHA-256 hash stored in MySQL."""

    def is_duplicate(self, file_bytes: bytes, user_id: str) -> Optional[Dict]:
        """Returns existing file doc if duplicate, None otherwise."""
        try:
            from botai.config.MySQL_config import get_db
            sha = hashlib.sha256(file_bytes).hexdigest()
            db = get_db()
            if db is None:
                return None
            existing = db.files.find_one({'user_id': user_id, 'sha256': sha})
            return existing
        except Exception as e:
            print(f"[DuplicateDetector] Error: {e}")
            return None


class ChunkManager:
    """Splits large text into overlapping chunks for RAG indexing."""

    def __init__(self, chunk_size: int = None, overlap: int = None):
        from botai.config import settings
        self.chunk_size = chunk_size or settings.RAG_CHUNK_SIZE
        self.overlap    = overlap    or settings.RAG_CHUNK_OVERLAP

    def chunk(self, text: str) -> List[Dict]:
        """Return a list of {text, chunk_index, start_char, end_char} dicts."""
        chunks = []
        start = 0
        idx = 0
        effective_overlap = min(self.overlap, self.chunk_size - 1) if self.overlap > 0 else 0
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            chunks.append({
                'text':        text[start:end],
                'chunk_index': idx,
                'start_char':  start,
                'end_char':    end
            })
            if end >= len(text):
                break
            advance = self.chunk_size - effective_overlap
            if advance <= 0:
                advance = 1
            start += advance
            idx += 1
        return chunks


# Global singletons
file_processor     = FileProcessor()
chunk_manager      = ChunkManager()
duplicate_detector = DuplicateDetector()
