"""
File Generator Service - Creates downloadable files from content
Supports: .docx, .pdf, .xlsx, .csv, .txt, .html, .md, .py, .json, .xml, .pptx
"""
import os
import re
import json
import csv
import io
import tempfile
from datetime import datetime

GENERATED_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'uploads', 'generated')
os.makedirs(GENERATED_DIR, exist_ok=True)


def generate_file(filename, content, file_type=None):
    """Generate a file and return the path. Content can be text or structured data."""
    if not file_type:
        file_type = filename.rsplit('.', 1)[-1].lower() if '.' in filename else 'txt'

    ext = file_type.lower()
    safe_name = re.sub(r'[^\w\-.]', '_', filename)
    safe_name = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{safe_name}"

    try:
        if ext in ('docx', 'doc'):
            path = _generate_docx(safe_name, content)
        elif ext == 'pdf':
            path = _generate_pdf(safe_name, content)
        elif ext in ('xlsx', 'xls'):
            path = _generate_xlsx(safe_name, content)
        elif ext == 'pptx':
            path = _generate_pptx(safe_name, content)
        elif ext == 'csv':
            path = _generate_csv(safe_name, content)
        elif ext == 'html':
            path = _generate_html_file(safe_name, content)
        elif ext == 'md':
            path = _generate_markdown(safe_name, content)
        elif ext == 'json':
            path = _generate_json_file(safe_name, content)
        elif ext == 'xml':
            path = _generate_xml_file(safe_name, content)
        else:
            path = _generate_text(safe_name, content)

        return {'success': True, 'path': path, 'filename': safe_name, 'size': os.path.getsize(path)}
    except Exception as e:
        print(f"[FileGen] Error generating {ext} file: {e}")
        return {'success': False, 'error': 'File generation failed'}


def _generate_docx(filename, content):
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    style = doc.styles['Normal']
    style.font.size = Pt(11)
    style.font.name = 'Calibri'

    lines = content.split('\n')
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('# '):
            p = doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith('## '):
            p = doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith('### '):
            p = doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith('- ') or stripped.startswith('* '):
            doc.add_paragraph(stripped[2:], style='List Bullet')
        elif stripped.startswith('---') or stripped.startswith('==='):
            doc.add_paragraph('_' * 50)
        elif stripped:
            doc.add_paragraph(stripped)
        else:
            doc.add_paragraph('')

    path = os.path.join(GENERATED_DIR, filename)
    doc.save(path)
    return path


def _generate_pdf(filename, content):
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font('Helvetica', size=11)

    lines = content.split('\n')
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('# '):
            pdf.set_font('Helvetica', 'B', 18)
            pdf.cell(0, 12, stripped[2:], new_x="LMARGIN", new_y="NEXT")
            pdf.set_font('Helvetica', size=11)
        elif stripped.startswith('## '):
            pdf.set_font('Helvetica', 'B', 14)
            pdf.cell(0, 10, stripped[3:], new_x="LMARGIN", new_y="NEXT")
            pdf.set_font('Helvetica', size=11)
        elif stripped.startswith('### '):
            pdf.set_font('Helvetica', 'B', 12)
            pdf.cell(0, 8, stripped[4:], new_x="LMARGIN", new_y="NEXT")
            pdf.set_font('Helvetica', size=11)
        elif stripped.startswith('- '):
            pdf.cell(5)
            pdf.cell(0, 7, f"  {stripped[2:]}", new_x="LMARGIN", new_y="NEXT")
        elif stripped:
            pdf.multi_cell(0, 7, stripped)
        else:
            pdf.ln(3)

    path = os.path.join(GENERATED_DIR, filename)
    pdf.output(path)
    return path


def _generate_xlsx(filename, content):
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    header_font = Font(bold=True, size=12, color="FFFFFF")
    header_fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
    border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    lines = content.strip().split('\n')
    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue
        # Try parsing as CSV-like data (pipe-separated or comma-separated)
        if '|' in stripped:
            cells = [c.strip() for c in stripped.split('|') if c.strip()]
        elif ',' in stripped:
            cells = [c.strip() for c in stripped.split(',') if c.strip()]
        else:
            cells = [stripped]

        for j, cell in enumerate(cells, 1):
            c = ws.cell(row=i+1, column=j, value=cell)
            c.border = border
            c.alignment = Alignment(wrap_text=True)
            if i == 0:
                c.font = header_font
                c.fill = header_fill

    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

    path = os.path.join(GENERATED_DIR, filename)
    wb.save(path)
    return path


def _generate_pptx(filename, content):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    lines = content.split('\n')
    title_text = "Presentation"
    bullet_points = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith('# '):
            title_text = stripped[2:]
        elif stripped.startswith('- ') or stripped.startswith('* '):
            bullet_points.append(stripped[2:])
        elif stripped and not bullet_points:
            title_text = stripped

    if title_text:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        if bullet_points and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            for bp in bullet_points[:10]:
                p = tf.add_paragraph()
                p.text = bp
                p.level = 0

        if not bullet_points:
            body = slide.placeholders[1]
            body.text = content[:500] if content else "Generated by CUTM AI"

    path = os.path.join(GENERATED_DIR, filename)
    prs.save(path)
    return path


def _generate_csv(filename, content):
    path = os.path.join(GENERATED_DIR, filename)
    lines = content.strip().split('\n')
    with open(path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            if '|' in stripped:
                writer.writerow([c.strip() for c in stripped.split('|') if c.strip()])
            elif ',' in stripped:
                writer.writerow([c.strip() for c in stripped.split(',') if c.strip()])
            else:
                writer.writerow([stripped])
    return path


def _generate_html_file(filename, content):
    path = os.path.join(GENERATED_DIR, filename)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    return path


def _generate_markdown(filename, content):
    path = os.path.join(GENERATED_DIR, filename)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    return path


def _generate_json_file(filename, content):
    path = os.path.join(GENERATED_DIR, filename)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    return path


def _generate_xml_file(filename, content):
    path = os.path.join(GENERATED_DIR, filename)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    return path


def _generate_text(filename, content):
    path = os.path.join(GENERATED_DIR, filename)
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    return path


def get_download_url(filename):
    return f"/api/files/download/{filename}"
